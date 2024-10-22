

import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches

# Giving Image path
img_path = '123.png'

# Creating an Presentation object
ppt = Presentation()

# Create a blank Slide
blank_slide_layout = ppt.slide_layouts[6]

slide = ppt.slides.add_slide(blank_slide_layout)

# For margins
# left = top = Inches(1)
# left = top = width = height = Inches(1)  
# using the add_textbox() method  
textBox = slide.shapes.add_textbox(100, 50, 1180, 27)  
# creating a text frame  
textFrame = textBox.text_frame  
  
# inserting text to the text frame  
textFrame.text = "link url :- https://www.instagram.com/p/CkzwKXbrA8e/"  

# adding images
pic = slide.shapes.add_picture(r"D:\Automation\sunset_400.jpeg",35,150)

left = Inches(2)
height = Inches(3)
top =Inches(2)
pic = slide.shapes.add_picture(r"D:\Automation\sunset_400.jpeg", left,top, height = height)
# # save file
ppt.save(r"D:\Automation\PPt\alsaya.pptx")

print("Done")


